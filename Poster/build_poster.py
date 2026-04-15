"""Generate the FYP poster and standee.

Two deliverables share content and theme but use different canvas aspect ratios:
  * FYP_Poster.pptx   — A1 portrait (23.4 x 33.1 in), based on Ref_1 layout.
  * FYP_Standee.pptx  — Roll-up banner 33 x 82 in, based on Ref_2 layout.

Theme: "Deep Plum + Sage" — see Report/poster_theme.md.
Layout: same section grid as Ref_1/Ref_2 — see Report/poster_layout.md.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
REPORT_ASSETS = os.path.join(os.path.dirname(HERE), "Report", "assets")
OUT_DIR = HERE

# =============================================================================
# Theme
# =============================================================================

PALETTE = {
    "primary":        RGBColor(0x2E, 0x1A, 0x47),   # deep plum
    "primary_accent": RGBColor(0x6B, 0x4E, 0x8C),   # mid plum
    "bright_accent":  RGBColor(0xF4, 0xA2, 0x61),   # warm amber
    "supporting":     RGBColor(0x86, 0xA8, 0x8E),   # sage
    "background":     RGBColor(0xFA, 0xF7, 0xF0),   # cream
    "card":           RGBColor(0xFF, 0xFF, 0xFF),
    "text":           RGBColor(0x1E, 0x1E, 0x2E),
    "text_muted":     RGBColor(0x5D, 0x5F, 0x76),
    "border":         RGBColor(0xE0, 0xDC, 0xD2),
}

FONT_TITLE = "Segoe UI Black"
FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_MONO = "Cascadia Mono Semibold"

# =============================================================================
# Project content
# =============================================================================

PROJECT_TITLE = "Smart Attendance System using Facial Image Recognition"
INSTITUTE = "GHULAM ISHAQ KHAN INSTITUTE OF ENGINEERING SCIENCES AND TECHNOLOGY"
FACULTY = "FACULTY OF ELECTRICAL ENGINEERING"
GIKI_LOGO = os.path.join(REPORT_ASSETS, "giki_logo.jpeg")

ABSTRACT = (
    "This project replaces manual classroom attendance with a single classroom "
    "photograph. A teacher uploads the image via a web dashboard; the system "
    "detects every face using MTCNN with a RetinaFace fallback, filters "
    "low-quality crops, embeds each face with FaceNet (512-d) followed by a "
    "learned 128-d projection head, and matches the embeddings against an "
    "enrollment store using minimum L2 distance with an SVM confirmation layer. "
    "The pipeline reaches 86% recognition on unseen classroom photos while "
    "processing each image in 1\u20135 seconds on commodity hardware."
)

OBJECTIVES = [
    "Eliminate roll call from classroom lectures.",
    "Reduce proxy / fraudulent attendance.",
    "Deliver auditable attendance records.",
    "Run on-premises with zero recurring cost.",
    "Preserve student privacy — no raw photos stored.",
]

APPLICATIONS = [
    "University lectures & tutorials.",
    "Examination hall attendance & identity verification.",
    "Event / workshop registration checks.",
    "Seminar & guest-lecture tracking.",
    "Faculty-level analytics on absenteeism trends.",
]

KEY_TOOLS = [
    "Python 3 \u2014 FastAPI, TensorFlow, OpenCV, scikit-learn",
    "FaceNet (keras_facenet), MTCNN, RetinaFace",
    "Custom projection head (triplet loss, 128-d)",
    "Node.js + Express + TypeScript + Prisma ORM",
    "React 18 + Vite + TypeScript (frontend)",
    "SQLite (development), PostgreSQL-ready (production)",
    "Nginx + Let's Encrypt + systemd (deployment)",
]

RESULT_STATS = [
    ("86%",   "Recognition rate"),
    ("243",   "Faces evaluated"),
    ("113",   "Students enrolled"),
    ("1\u20135 s",  "Per-photo latency"),
]

RESULT_NOTE = (
    "Evaluated on 14 unseen classroom photographs drawn from the 2022 and 2023 "
    "batches. Nine of fourteen photos reached 100% recognition; remaining errors "
    "cluster on small / pose-extreme faces in dense groups, which the teacher-"
    "review UI flags for manual approval."
)

CONCLUSION = (
    "The Smart Attendance System demonstrates that a single classroom "
    "photograph, processed by an on-premises deep-learning pipeline, is a "
    "practical replacement for manual roll call. The design is privacy-"
    "respecting (no raw images stored), cost-free to operate (open-source stack "
    "on a free-tier VM), and ships with a teacher-review step that guards "
    "against residual recognition errors."
)

GROUP_MEMBERS = [
    ("Hassaan Ahmed", "FCSE"),
    ("Osaid Khan Afridi", "FCSE"),
    ("Saad Khan", "FEE"),
    ("Najaf", "FEE"),
]

ADVISOR = ("Dr. Zaiwar Ali", "Advisor", "Faculty of Electrical Engineering")
CO_ADVISOR = ("Dr. Arbab Abdur Rahim", "Co-Advisor / Dean",
              "Faculty of Electrical Engineering")

FLOWCHART_IMG = os.path.join(REPORT_ASSETS, "fig_pipeline.png")
SYSTEM_IMG = os.path.join(REPORT_ASSETS, "fig_architecture.png")
RESULTS_IMG = os.path.join(REPORT_ASSETS, "fig_accuracy.png")


# =============================================================================
# Helpers
# =============================================================================

def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width=Pt(0.75)):
    line = shape.line
    line.color.rgb = rgb
    line.width = width


def no_line(shape):
    line = shape.line
    line.fill.background()


def add_rect(slide, x, y, w, h, fill, line=None, corner=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    set_fill(shp, fill)
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line)
    if corner > 0:
        # corner is 0..100 (%)
        shp.adjustments[0] = corner / 100.0
    return shp


def add_text(slide, x, y, w, h, text,
             font_name=FONT_BODY, font_size=11, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = Inches(0.05)
    tb.text_frame.margin_right = Inches(0.05)
    tb.text_frame.margin_top = Inches(0.02)
    tb.text_frame.margin_bottom = Inches(0.02)
    tb.text_frame.vertical_anchor = anchor
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def add_multiline_text(slide, x, y, w, h, lines,
                       font_name=FONT_BODY, font_size=11, color=None,
                       bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                       line_spacing=1.2):
    """lines = list[str]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return tb


def add_section_header(slide, x, y, w, text, h=None, font_size=18):
    """Solid plum pill with uppercase white title text."""
    if h is None:
        h = Inches(0.55)
    pill = add_rect(slide, x, y, w, h, PALETTE["primary"], corner=18)
    # text inside
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + text.upper()
    run.font.name = FONT_HEAD
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = PALETTE["card"]
    return pill


def add_bullets_sage(slide, x, y, w, h, items, font_size=11):
    """Bullet list with small sage triangles prefixed."""
    lines = ["\u25B8 " + item for item in items]
    return add_multiline_text(slide, x, y, w, h, lines,
                              font_name=FONT_BODY, font_size=font_size,
                              color=PALETTE["text"], line_spacing=1.25)


def add_card(slide, x, y, w, h):
    return add_rect(slide, x, y, w, h, PALETTE["card"],
                    line=PALETTE["border"], corner=3)


def add_image_fitted(slide, path, x, y, w, h):
    """Add an image, scaled to fit inside the given rect without distortion."""
    if not os.path.exists(path):
        return None
    from PIL import Image as PILImage
    im = PILImage.open(path)
    iw, ih = im.size
    img_aspect = iw / ih
    box_aspect = w / h
    if img_aspect > box_aspect:
        # fit to width
        new_w = w
        new_h = int(w / img_aspect)
        new_x = x
        new_y = y + (h - new_h) // 2
    else:
        new_h = h
        new_w = int(h * img_aspect)
        new_y = y
        new_x = x + (w - new_w) // 2
    return slide.shapes.add_picture(path, new_x, new_y, new_w, new_h)


# =============================================================================
# Shared banner builders
# =============================================================================

def build_header_strip(slide, x, y, w, h):
    """Institute name + faculty subtitle + GIKI logo on the right."""
    add_rect(slide, x, y, w, h, PALETTE["background"])
    # logo on right
    if os.path.exists(GIKI_LOGO):
        logo_side = h - Inches(0.1)
        logo_x = x + w - logo_side - Inches(0.15)
        logo_y = y + (h - logo_side) // 2
        slide.shapes.add_picture(GIKI_LOGO, logo_x, logo_y, logo_side, logo_side)
    # text region left of logo
    text_w = w - Inches(1.6)
    add_text(slide, x + Inches(0.2), y, text_w, h // 2, INSTITUTE,
             font_name=FONT_HEAD, font_size=14, bold=True,
             color=PALETTE["primary"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, x + Inches(0.2), y + h // 2, text_w, h // 2, FACULTY,
             font_name=FONT_HEAD, font_size=12, bold=True,
             color=PALETTE["primary_accent"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)


def build_title_band(slide, x, y, w, h, size=36):
    add_rect(slide, x, y, w, h, PALETTE["primary"], corner=2)
    # Force two-line title for visual punch (it also wraps naturally in narrow widths)
    add_multiline_text(slide, x, y, w, h,
                       ["SMART ATTENDANCE SYSTEM",
                        "USING FACIAL IMAGE RECOGNITION"],
                       font_name=FONT_TITLE, font_size=size, bold=True,
                       color=PALETTE["card"], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def build_results_stats(slide, x, y, w, h):
    """4-up grid of big stat numbers + small labels."""
    add_card(slide, x, y, w, h)
    pad = Inches(0.1)
    inner_w = w - 2 * pad
    inner_h = h - 2 * pad
    cell_w = inner_w // 2
    cell_h = inner_h // 2
    for i, (num, label) in enumerate(RESULT_STATS):
        col = i % 2
        row = i // 2
        cx = x + pad + col * cell_w
        cy = y + pad + row * cell_h
        add_text(slide, cx, cy, cell_w, int(cell_h * 0.65), num,
                 font_name=FONT_MONO, font_size=28, bold=True,
                 color=PALETTE["bright_accent"], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.BOTTOM)
        add_text(slide, cx, cy + int(cell_h * 0.60), cell_w, int(cell_h * 0.4), label,
                 font_name=FONT_HEAD, font_size=11, bold=True,
                 color=PALETTE["primary"], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP)


def build_members_advisors(slide, x, y, w, h, orient="poster"):
    """Two side-by-side bands: members on left, advisors on right."""
    # Left band (members)
    left_w = int(w * 0.60)
    right_w = w - left_w - Inches(0.15)
    right_x = x + left_w + Inches(0.15)

    add_rect(slide, x, y, left_w, h, PALETTE["primary"], corner=2)
    # label
    add_text(slide, x + Inches(0.2), y + Inches(0.1),
             left_w - Inches(0.4), Inches(0.4), "GROUP MEMBERS",
             font_name=FONT_HEAD, font_size=14, bold=True,
             color=PALETTE["bright_accent"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    # names
    cols = 2
    rows = 2
    name_start_y = y + Inches(0.55)
    name_w = (left_w - Inches(0.4)) // cols
    name_h = (h - Inches(0.7)) // rows
    for i, (name, dept) in enumerate(GROUP_MEMBERS):
        col = i % cols
        row = i // cols
        nx = x + Inches(0.2) + col * name_w
        ny = name_start_y + row * name_h
        add_text(slide, nx, ny, name_w, Pt(16), name,
                 font_name=FONT_HEAD, font_size=12, bold=True,
                 color=PALETTE["card"], align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP)
        add_text(slide, nx, ny + Inches(0.25), name_w, Pt(12), dept,
                 font_name=FONT_BODY, font_size=10,
                 color=PALETTE["supporting"], align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP)

    # Right band (advisors)
    add_rect(slide, right_x, y, right_w, h,
             PALETTE["primary_accent"], corner=2)
    adv_lines_y = y + Inches(0.1)
    # Advisor
    add_text(slide, right_x + Inches(0.2), adv_lines_y,
             right_w - Inches(0.4), Inches(0.3), "ADVISOR",
             font_name=FONT_HEAD, font_size=11, bold=True,
             color=PALETTE["bright_accent"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    add_text(slide, right_x + Inches(0.2), adv_lines_y + Inches(0.30),
             right_w - Inches(0.4), Inches(0.3), ADVISOR[0],
             font_name=FONT_HEAD, font_size=12, bold=True,
             color=PALETTE["card"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    # Co-Advisor
    co_y = adv_lines_y + Inches(0.80)
    add_text(slide, right_x + Inches(0.2), co_y,
             right_w - Inches(0.4), Inches(0.3), "CO-ADVISOR",
             font_name=FONT_HEAD, font_size=11, bold=True,
             color=PALETTE["bright_accent"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    add_text(slide, right_x + Inches(0.2), co_y + Inches(0.30),
             right_w - Inches(0.4), Inches(0.3), CO_ADVISOR[0],
             font_name=FONT_HEAD, font_size=12, bold=True,
             color=PALETTE["card"], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)


def build_partners_strip(slide, x, y, w, h):
    add_rect(slide, x, y, w, h, PALETTE["supporting"], corner=2)
    add_text(slide, x, y, w, h,
             "Built at GIK Institute  \u2022  Faculty of Electrical Engineering  "
             "\u2022  Senior Design Project 2025-26",
             font_name=FONT_HEAD, font_size=12, bold=True,
             color=PALETTE["card"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
# POSTER (Ref_1 layout — portrait A1-ish)
# =============================================================================

def build_poster():
    prs = Presentation()
    # A1 portrait, 23.4 x 33.1 inches
    W = Inches(23.4)
    H = Inches(33.1)
    prs.slide_width = W
    prs.slide_height = H

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, W, H, PALETTE["background"])

    MARGIN = Inches(0.5)
    GUT = Inches(0.3)
    inner_w = W - 2 * MARGIN

    # --- Header strip ---
    hs_h = Inches(1.2)
    build_header_strip(slide, MARGIN, MARGIN, inner_w, hs_h)

    # --- Title band ---
    tb_y = MARGIN + hs_h + Inches(0.1)
    tb_h = Inches(1.6)
    build_title_band(slide, MARGIN, tb_y, inner_w, tb_h, size=36)

    # --- Content grid ---
    content_y = tb_y + tb_h + Inches(0.25)
    # Budget from content_y to end of slide
    remaining = H - content_y - MARGIN
    # Reserve space for partners strip + members block at bottom
    bottom_reserved = Inches(0.55) + Inches(0.2) + Inches(1.6)
    grid_h = remaining - bottom_reserved - Inches(0.25)

    left_w = int(inner_w * 0.42)
    right_w = inner_w - left_w - GUT
    left_x = MARGIN
    right_x = MARGIN + left_w + GUT

    # Split grid into 4 rows
    row_gap = Inches(0.25)
    total_gaps = 3 * row_gap
    usable_rows = grid_h - total_gaps
    # Proportions: 26% / 17% / 26% / 31%
    row1_h = int(usable_rows * 0.28)
    row2_h = int(usable_rows * 0.19)
    row3_h = int(usable_rows * 0.25)
    row4_h = usable_rows - row1_h - row2_h - row3_h

    hdr_h = Inches(0.55)

    # Row 1: Abstract (L) + Flowchart (R)
    y = content_y
    add_section_header(slide, left_x, y, left_w, "Abstract",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.08)
    body_h = row1_h - hdr_h - Inches(0.08)
    add_card(slide, left_x, body_y, left_w, body_h)
    add_text(slide, left_x + Inches(0.3), body_y + Inches(0.15),
             left_w - Inches(0.6), body_h - Inches(0.3),
             ABSTRACT,
             font_name=FONT_BODY, font_size=13, color=PALETTE["text"],
             align=PP_ALIGN.JUSTIFY, line_spacing=1.3)

    add_section_header(slide, right_x, y, right_w, "Recognition Pipeline",
                       h=hdr_h, font_size=18)
    add_card(slide, right_x, body_y, right_w, body_h)
    add_image_fitted(slide, FLOWCHART_IMG,
                     right_x + Inches(0.2), body_y + Inches(0.15),
                     right_w - Inches(0.4), body_h - Inches(0.3))

    # Row 2: Objectives (L) + Applications (R)
    y = content_y + row1_h + row_gap
    add_section_header(slide, left_x, y, left_w, "Objectives",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.08)
    body_h = row2_h - hdr_h - Inches(0.08)
    add_card(slide, left_x, body_y, left_w, body_h)
    add_bullets_sage(slide, left_x + Inches(0.3), body_y + Inches(0.15),
                     left_w - Inches(0.6), body_h - Inches(0.3),
                     OBJECTIVES, font_size=12)
    add_section_header(slide, right_x, y, right_w, "Applications",
                       h=hdr_h, font_size=18)
    add_card(slide, right_x, body_y, right_w, body_h)
    add_bullets_sage(slide, right_x + Inches(0.3), body_y + Inches(0.15),
                     right_w - Inches(0.6), body_h - Inches(0.3),
                     APPLICATIONS, font_size=12)

    # Row 3: Key Tools (L) + Results (R)
    y = content_y + row1_h + row2_h + 2 * row_gap
    add_section_header(slide, left_x, y, left_w, "Key Tools",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.08)
    body_h = row3_h - hdr_h - Inches(0.08)
    add_card(slide, left_x, body_y, left_w, body_h)
    add_bullets_sage(slide, left_x + Inches(0.3), body_y + Inches(0.15),
                     left_w - Inches(0.6), body_h - Inches(0.3),
                     KEY_TOOLS, font_size=11)

    add_section_header(slide, right_x, y, right_w, "Results",
                       h=hdr_h, font_size=18)
    stats_h = int(body_h * 0.55)
    build_results_stats(slide, right_x, body_y, right_w, stats_h)
    note_y = body_y + stats_h + Inches(0.12)
    note_h = body_h - stats_h - Inches(0.12)
    add_card(slide, right_x, note_y, right_w, note_h)
    add_text(slide, right_x + Inches(0.25), note_y + Inches(0.12),
             right_w - Inches(0.5), note_h - Inches(0.24),
             RESULT_NOTE,
             font_name=FONT_BODY, font_size=11, color=PALETTE["text"],
             align=PP_ALIGN.JUSTIFY, line_spacing=1.3)

    # Row 4: System (L) + Conclusion (R)
    y = content_y + row1_h + row2_h + row3_h + 3 * row_gap
    add_section_header(slide, left_x, y, left_w, "System Architecture",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.08)
    body_h = row4_h - hdr_h - Inches(0.08)
    add_card(slide, left_x, body_y, left_w, body_h)
    add_image_fitted(slide, SYSTEM_IMG,
                     left_x + Inches(0.2), body_y + Inches(0.15),
                     left_w - Inches(0.4), body_h - Inches(0.3))

    add_section_header(slide, right_x, y, right_w, "Conclusion",
                       h=hdr_h, font_size=18)
    add_card(slide, right_x, body_y, right_w, body_h)
    add_text(slide, right_x + Inches(0.3), body_y + Inches(0.2),
             right_w - Inches(0.6), body_h - Inches(0.4),
             CONCLUSION,
             font_name=FONT_BODY, font_size=13, color=PALETTE["text"],
             align=PP_ALIGN.JUSTIFY, line_spacing=1.35)

    # Partners strip
    y = content_y + grid_h + Inches(0.15)
    partners_h = Inches(0.55)
    build_partners_strip(slide, MARGIN, y, inner_w, partners_h)

    # Members / advisors
    ma_y = y + partners_h + Inches(0.2)
    ma_h = H - ma_y - MARGIN
    build_members_advisors(slide, MARGIN, ma_y, inner_w, ma_h)

    out_path = os.path.join(OUT_DIR, "FYP_Poster.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# =============================================================================
# STANDEE (Ref_2 layout — tall roll-up banner)
# =============================================================================

def build_standee():
    prs = Presentation()
    # 22 x 55 inches (~0.40 aspect, same as Ref_2)
    W = Inches(22.0)
    H = Inches(55.0)
    prs.slide_width = W
    prs.slide_height = H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, PALETTE["background"])

    MARGIN = Inches(0.5)
    GUT = Inches(0.3)
    inner_w = W - 2 * MARGIN
    col_w = (inner_w - GUT) // 2
    x_left = MARGIN
    x_right = MARGIN + col_w + GUT

    hdr_h = Inches(0.65)

    y = MARGIN

    # Header strip
    hs_h = Inches(1.15)
    build_header_strip(slide, MARGIN, y, inner_w, hs_h)
    y += hs_h + Inches(0.15)

    # Title band — kept narrower so title wraps cleanly on the tall standee
    tb_h = Inches(3.2)
    build_title_band(slide, MARGIN, y, inner_w, tb_h, size=24)
    y += tb_h + Inches(0.3)

    # Abstract full width
    add_section_header(slide, MARGIN, y, inner_w, "Abstract",
                       h=hdr_h, font_size=20)
    y += hdr_h + Inches(0.1)
    abs_h = Inches(4.8)
    add_card(slide, MARGIN, y, inner_w, abs_h)
    add_text(slide, MARGIN + Inches(0.35), y + Inches(0.25),
             inner_w - Inches(0.7), abs_h - Inches(0.5),
             ABSTRACT, font_name=FONT_BODY, font_size=14,
             color=PALETTE["text"], align=PP_ALIGN.JUSTIFY, line_spacing=1.4)
    y += abs_h + Inches(0.4)

    # Objectives + Flowchart
    row_h = Inches(6.2)
    add_section_header(slide, x_left, y, col_w, "Objectives",
                       h=hdr_h, font_size=18)
    add_section_header(slide, x_right, y, col_w, "Pipeline",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.1)
    body_h = row_h - hdr_h - Inches(0.1)
    add_card(slide, x_left, body_y, col_w, body_h)
    add_bullets_sage(slide, x_left + Inches(0.3), body_y + Inches(0.25),
                     col_w - Inches(0.6), body_h - Inches(0.5),
                     OBJECTIVES, font_size=13)
    add_card(slide, x_right, body_y, col_w, body_h)
    add_image_fitted(slide, FLOWCHART_IMG,
                     x_right + Inches(0.2), body_y + Inches(0.2),
                     col_w - Inches(0.4), body_h - Inches(0.4))
    y += row_h + Inches(0.4)

    # Applications + Key Tools
    row_h = Inches(5.6)
    add_section_header(slide, x_left, y, col_w, "Applications",
                       h=hdr_h, font_size=18)
    add_section_header(slide, x_right, y, col_w, "Key Tools",
                       h=hdr_h, font_size=18)
    body_y = y + hdr_h + Inches(0.1)
    body_h = row_h - hdr_h - Inches(0.1)
    add_card(slide, x_left, body_y, col_w, body_h)
    add_bullets_sage(slide, x_left + Inches(0.3), body_y + Inches(0.25),
                     col_w - Inches(0.6), body_h - Inches(0.5),
                     APPLICATIONS, font_size=13)
    add_card(slide, x_right, body_y, col_w, body_h)
    add_bullets_sage(slide, x_right + Inches(0.3), body_y + Inches(0.25),
                     col_w - Inches(0.6), body_h - Inches(0.5),
                     KEY_TOOLS, font_size=11)
    y += row_h + Inches(0.4)

    # System architecture full width
    add_section_header(slide, MARGIN, y, inner_w, "System Architecture",
                       h=hdr_h, font_size=20)
    y += hdr_h + Inches(0.1)
    sys_h = Inches(6.5)
    add_card(slide, MARGIN, y, inner_w, sys_h)
    add_image_fitted(slide, SYSTEM_IMG,
                     MARGIN + Inches(0.3), y + Inches(0.3),
                     inner_w - Inches(0.6), sys_h - Inches(0.6))
    y += sys_h + Inches(0.4)

    # Results (stats + chart side by side)
    add_section_header(slide, MARGIN, y, inner_w, "Results",
                       h=hdr_h, font_size=20)
    y += hdr_h + Inches(0.1)
    res_h = Inches(6.2)
    stats_w = int(inner_w * 0.45)
    chart_x = MARGIN + stats_w + Inches(0.3)
    chart_w = inner_w - stats_w - Inches(0.3)
    build_results_stats(slide, MARGIN, y, stats_w, res_h)
    add_card(slide, chart_x, y, chart_w, res_h)
    add_image_fitted(slide, RESULTS_IMG,
                     chart_x + Inches(0.2), y + Inches(0.2),
                     chart_w - Inches(0.4), res_h - Inches(0.4))
    y += res_h + Inches(0.4)

    # Conclusion
    add_section_header(slide, MARGIN, y, inner_w, "Conclusion",
                       h=hdr_h, font_size=20)
    y += hdr_h + Inches(0.1)
    concl_h = Inches(3.4)
    add_card(slide, MARGIN, y, inner_w, concl_h)
    add_text(slide, MARGIN + Inches(0.35), y + Inches(0.3),
             inner_w - Inches(0.7), concl_h - Inches(0.6),
             CONCLUSION, font_name=FONT_BODY, font_size=14,
             color=PALETTE["text"], align=PP_ALIGN.JUSTIFY, line_spacing=1.4)
    y += concl_h + Inches(0.4)

    # Partners strip
    partners_h = Inches(0.65)
    build_partners_strip(slide, MARGIN, y, inner_w, partners_h)
    y += partners_h + Inches(0.25)

    # Members + advisors
    ma_h = H - MARGIN - y
    build_members_advisors(slide, MARGIN, y, inner_w, ma_h, orient="standee")

    out_path = os.path.join(OUT_DIR, "FYP_Standee.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# =============================================================================
# Render PNG previews using python-pptx + PIL (simplified — uses first slide)
# =============================================================================

def render_preview_via_image_composition(pptx_path, png_path, width_px=1500):
    """
    Since LibreOffice is not installed, compose a simple PNG preview by
    re-drawing the slide contents with PIL. This is a low-fi preview only —
    the .pptx itself is the deliverable.

    For brevity, we only render the final .pptx dimensions as a placeholder
    image showing "See .pptx for full design". A full pixel-exact render would
    require LibreOffice or python-pptx-to-image tooling.
    """
    # Skip for now; the pptx file is the deliverable.
    pass


# =============================================================================
# Main
# =============================================================================

if __name__ == "__main__":
    os.makedirs(OUT_DIR, exist_ok=True)
    build_poster()
    build_standee()
    print("\nDone. Open the .pptx files in PowerPoint to preview.")
